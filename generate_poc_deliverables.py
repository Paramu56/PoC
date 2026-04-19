"""
One-off generator: Word write-up + two ~15-slide PowerPoint decks (executive + technical).
Run: python generate_poc_deliverables.py
Outputs: deliverables/*.docx, deliverables/*.pptx
"""

from __future__ import annotations

import os
from datetime import date

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt

OUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "deliverables")
TODAY = date.today().isoformat()


def _add_heading(doc: Document, text: str, level: int = 1) -> None:
    doc.add_heading(text, level=level)


def _p(doc: Document, text: str) -> None:
    doc.add_paragraph(text)


def build_docx(path: str) -> None:
    doc = Document()
    title = doc.add_heading("Karnataka Government Schemes — AI Proof of Concept", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run(f"Architecture, modules, technology, and roadmap\n{TODAY}")
    r.italic = True
    doc.add_paragraph()

    sections = [
        (
            "1. Purpose of this document",
            (
                "This document describes a local-first knowledge system built as a proof of concept (PoC) "
                "for Karnataka government schemes. It covers ingestion, retrieval-augmented generation (RAG), "
                "orchestration with Google Gemini, user interfaces (Streamlit and Flask), PDF kit generation, "
                "and optional location-aware features. It also records planned but not yet implemented "
                "capabilities—guardrails, privacy controls, and governance—so stakeholders can see both "
                "what exists today and what a production-oriented phase would add.",
            ),
        ),
        (
            "2. What the PoC demonstrates",
            (
                "The PoC shows an end-to-end AI application—not only a chat interface. It includes: "
                "(1) ingestion of scheme text from PDFs into a searchable vector index with rich metadata; "
                "(2) semantic search over that index; (3) an orchestrated RAG pipeline that retrieves context, "
                "optionally refines queries, augments with rule-like graph context, and asks an LLM for a "
                "grounded answer with citation-oriented context blocks; (4) a Streamlit demo with profile-aware "
                "questions, scheme ranking, and downloadable application-kit PDFs; (5) a small Flask API "
                "exposing the same RAG core; (6) optional geocoding and nearest Bangalore One centre suggestion. "
                "For PDF kits, the system can prefer full scheme text extracted directly from the official "
                "source PDF (numbered sections), falling back to indexed chunks when needed.",
            ),
        ),
        (
            "3. High-level architecture",
            (
                "Data flows: Karnataka Schemes.pdf → ingestion script → Chroma persistent vector store "
                "(sentence-transformer embeddings). At query time, the orchestration module queries Chroma, "
                "may run an optional Gemini step to judge sufficiency and suggest refined search strings, "
                "merges and deduplicates chunks, appends filtered statements from a compiled graph JSON "
                "(derived from Cypher-style rules, without Neo4j), and sends a bounded context block to Gemini "
                "for the final answer. Streamlit and Flask call the same orchestrated function and receive "
                "answers, citations metadata, and debug information.",
            ),
        ),
        (
            "4. Main modules and responsibilities",
            (
                "ingest_karnataka_schemes.py — Extracts text with pypdf, parses scheme blocks heuristically, "
                "chunks and stores in Chroma with metadata (scheme name, page, category, tags, source). "
                "Can compile schemes.cypher into graph_compiled.json.\n\n"
                "orchestrated_rag_schemes.py — Core RAG: Chroma retrieval, context formatting, optional "
                "Gemini-driven refinement and second-pass retrieval, deduplication, graph add-on, final Gemini "
                "answer; exposes run_orchestrated_rag_result for applications.\n\n"
                "graph_knowledge.py — Loads and filters Cypher-like rules for schemes relevant to retrieved chunks.\n\n"
                "graph_scheme_documents.py — Maps schemes to required documents via DEPENDS_ON edges parsed from the Cypher file.\n\n"
                "user_profile.py — Builds structured questions from profile fields (shared by Flask and Streamlit).\n\n"
                "scheme_choice_label.py — Produces readable scheme labels from retrieval rows.\n\n"
                "scheme_source_pdf.py — Splits the official PDF on numbered scheme headings; returns the best-matching section for kit text.\n\n"
                "poc_pdf.py — Builds application-kit PDFs with fpdf2 (text normalization and reflow).\n\n"
                "poc_nearby.py — Nominatim geocoding and nearest Bangalore One centre from CSV.\n\n"
                "poc_streamlit.py — Primary PoC UI: orchestrated RAG, picker, PDF download, source-PDF captions.\n\n"
                "app.py — Flask UI and /api/ask JSON endpoint using the same pipeline.\n\n"
                "query_schemes.py — CLI query utility. Additional utilities include compare_schemes_pdf_vs_db.py and list_uploaded_schemes.py.",
            ),
        ),
        (
            "5. Technology stack",
            (
                "Language: Python 3.\n"
                "Vector database: Chroma (persistent on disk).\n"
                "Embeddings: sentence-transformers via Chroma (default: all-MiniLM-L6-v2).\n"
                "PDF reading: pypdf.\n"
                "LLM: Google Gemini (google-genai SDK; e.g. gemini-flash-latest / gemini-pro-latest).\n"
                "Demos: Streamlit; alternate web stack: Flask with templates and static assets.\n"
                "PDF generation: fpdf2.\n"
                "Geocoding (PoC): OpenStreetMap Nominatim (HTTP; policy-compliant usage).\n"
                "PoC dependency file: requirements-poc.txt lists Streamlit, fpdf2, google-genai; Chroma, pypdf, "
                "and sentence-transformers are part of the broader environment used with ingest and RAG.",
            ),
        ),
        (
            "6. Orchestration flow (RAG core)",
            (
                "(1) Initial retrieval of top-k chunks from Chroma.\n"
                "(2) Optional Gemini call: assess whether context is sufficient; if not, suggest refined queries "
                "(configurable; default may skip refinement to save quota).\n"
                "(3) Optional additional retrievals; deduplicate by scheme, page, and text.\n"
                "(4) Truncate chunks per GEMINI_CONTEXT_CHUNK_CHARS to respect API limits.\n"
                "(5) Append graph-derived rules from graph_compiled.json when relevant.\n"
                "(6) Final Gemini generation with numbered context lines for traceability.\n"
                "(7) Return answer, citations-style chunk references, and debug fields (retrieval stats, model used, etc.).",
            ),
        ),
        (
            "7. Alternative technologies (for scaling or customer environments)",
            (
                "Vector DB: Qdrant, Weaviate, Milvus, or pgvector on PostgreSQL.\n"
                "Embeddings: hosted APIs (OpenAI, Cohere, Vertex) or larger open models (e.g. bge-m3) where latency and cost allow.\n"
                "LLM: OpenAI, Anthropic, Azure OpenAI, or self-hosted open weights via vLLM/Ollama.\n"
                "Orchestration frameworks: LangChain, LlamaIndex, Haystack, Semantic Kernel—this PoC uses explicit Python for clarity and control.\n"
                "UI/API: FastAPI, Next.js, Gradio.\n"
                "Graph: Neo4j or other property graphs if rules outgrow file-based compilation.\n"
                "PDF: pdfplumber, Unstructured, or cloud OCR for scanned documents.\n"
                "Trade-offs involve cost, latency, compliance, and team skills—not a single universal best choice.",
            ),
        ),
        (
            "8. Reference framing for AI programs (governance narrative)",
            (
                "For enterprise conversations, the solution can be described in layers: "
                "(1) Data and ingestion, (2) Knowledge storage and rules, (3) Orchestration and grounding, "
                "(4) Generation and transparency, (5) User experience and exports, (6) Operations and governance. "
                "Formal frameworks such as the NIST AI Risk Management Framework (Govern, Map, Measure, Manage) "
                "or ISO/IEC 42001 (AI management systems) can align roadmap discussions with customer risk and compliance teams. "
                "The EU AI Act may matter for customers with EU deployment or high-risk use cases.",
            ),
        ),
        (
            "9. Planned: guardrails and safe behaviour (not yet fully implemented)",
            (
                "The following are design directions discussed for a subsequent phase; they reduce hallucination risk "
                "and unsafe outputs beyond what raw RAG provides.\n\n"
                "Retrieval: minimum relevance thresholds on vector distance; require minimum chunk count or text mass; "
                "scheme-alignment checks when the question names a specific programme.\n\n"
                "Generation: strict prompts requiring answers to stick to provided context; structured outputs "
                "(e.g. answer + confidence + citations); optional second-pass verification that flags unsupported claims.\n\n"
                "Abstention: when retrieval is weak, return a clear “insufficient context” response instead of inventing detail.\n\n"
                "Output: basic filters for disallowed content; optional PII detection/redaction in logs and responses.\n\n"
                "Product: visible source snippets or citation lists beside answers; stronger disclaimers pointing users to official portals.",
            ),
        ),
        (
            "10. Planned: privacy and governance (not yet fully implemented)",
            (
                "Privacy: classify what user data may be sent to cloud LLM providers; minimize and avoid unnecessary PII; "
                "short or no retention of prompts in logs for PoC extensions; secure handling of API keys; "
                "awareness that the Streamlit “network URL” exposes the demo to other devices on the LAN.\n\n"
                "Governance: role-based access to repos and keys; separation of dev vs demo API quotas; "
                "versioning of ingested corpora and prompts; audit metadata (timestamps, model ID, index version) "
                "without storing full transcripts where not required; alignment with organizational AI use policies "
                "and vendor terms (e.g. Google Gemini data processing).\n\n"
                "These items support customer due diligence and can be prioritized per contract or sector.",
            ),
        ),
        (
            "11. Current strengths and honest limits",
            (
                "Strengths: End-to-end pipeline; local index; citation-oriented context formatting; orchestrated multi-step RAG; "
                "dual interfaces; PDF export; graph-flavoured rules without a heavy database; source-PDF sections for kit fidelity.\n\n"
                "Limits: Guardrails and formal privacy/governance controls are roadmap items, not complete in the PoC. "
                "Production would add hardening, monitoring, SLAs, and integration with customer identity and data policies.",
            ),
        ),
        (
            "12. Conclusion",
            (
                "This PoC demonstrates that a credible AI system can be built around grounded retrieval, transparent context blocks, "
                "and practical user workflows—including document export—while leaving a clear path to stronger safety, privacy, "
                "and governance for customer-specific deployments.",
            ),
        ),
    ]

    for heading, body in sections:
        text = body[0] if isinstance(body, tuple) else body
        _add_heading(doc, heading, level=1)
        for para in text.split("\n\n"):
            _p(doc, para.strip())

    doc.save(path)
    print(f"Wrote {path}")


def _slide_title_only(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    tx = slide.shapes.add_textbox(PptInches(0.5), PptInches(2.2), PptInches(9), PptInches(1.2))
    tf = tx.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = PptPt(32)
    tf.paragraphs[0].font.bold = True
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = PptPt(14)
        p.space_before = PptPt(12)


def _slide_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide_layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = PptPt(18)


def build_ppt_executive(path: str) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)

    _slide_title_only(
        prs,
        "Karnataka Schemes AI — Proof of Concept",
        f"Executive overview · {TODAY}",
    )
    _slide_bullets(
        prs,
        "What we built",
        [
            "An end-to-end assistant over government scheme content—not just a chat box",
            "Grounded answers using your indexed documents plus optional rule context",
            "Working demo: questions, scheme shortlist, downloadable application-kit PDF",
        ],
    )
    _slide_bullets(
        prs,
        "Business value",
        [
            "Faster discovery of relevant schemes for citizens and field staff",
            "Consistent packaging of scheme text into a printable kit (PoC)",
            "Same engine usable from a web UI or API for future products",
        ],
    )
    _slide_bullets(
        prs,
        "How it works (one minute)",
        [
            "Official PDFs are ingested into a private searchable knowledge base",
            "User questions pull the most relevant passages first, then an AI model words the answer",
            "Optional steps improve search quality before the final answer",
        ],
    )
    _slide_bullets(
        prs,
        "Why this architecture matters",
        [
            "Retrieval-first design ties answers to source material",
            "Orchestration allows refinement and transparency (what was read)",
            "Export to PDF supports real-world “takeaway” workflows",
        ],
    )
    _slide_bullets(
        prs,
        "Technology (at a glance)",
        [
            "Python application; local vector database (Chroma)",
            "Google Gemini for language generation",
            "Streamlit demo + Flask API option",
        ],
    )
    _slide_bullets(
        prs,
        "Flexibility for your environment",
        [
            "Vector DB, embeddings, and LLM can be swapped for cloud or on-prem options",
            "Orchestration can later use standard frameworks if your IT prefers",
            "We document trade-offs: cost, latency, compliance—not one-size-fits-all",
        ],
    )
    _slide_bullets(
        prs,
        "Roadmap: trust and safety (planned)",
        [
            "Stronger grounding: relevance thresholds and “we don’t know” when data is thin",
            "Clear citations and optional verification of claims against sources",
            "Disclaimers: official portals remain authoritative",
        ],
    )
    _slide_bullets(
        prs,
        "Roadmap: privacy and governance (planned)",
        [
            "Data minimization and policies for what may be sent to AI providers",
            "Audit trail of model and knowledge-base version—not just free-text logs",
            "Access control and separation of dev vs production keys and quotas",
        ],
    )
    _slide_bullets(
        prs,
        "Alignment with enterprise AI practice",
        [
            "Layered story: data → knowledge → orchestration → experience → governance",
            "Can map to NIST AI RMF or ISO/IEC 42001 conversations when you need it",
        ],
    )
    _slide_bullets(
        prs,
        "What is proven vs what comes next",
        [
            "Proven: working pipeline, demo UX, API hook, PDF export",
            "Next: production hardening, guardrails, privacy controls, SLAs—scoped to your sector",
        ],
    )
    _slide_bullets(
        prs,
        "Suggested next conversation",
        [
            "Pilot scope: corpus, users, compliance constraints",
            "Hosting: your cloud, VPC, or air-gapped options",
            "Success metrics: accuracy checks, user satisfaction, time saved",
        ],
    )
    _slide_bullets(
        prs,
        "Who this serves",
        [
            "Citizens and field staff exploring eligibility and steps",
            "Program teams needing consistent scheme text in one place",
            "Partners integrating via API into existing portals or contact centres",
        ],
    )
    _slide_bullets(
        prs,
        "Ways to engage",
        [
            "Pilot: limited corpus, defined users, success criteria",
            "Co-development: guardrails and hosting aligned to your IT policy",
            "Support: iteration on accuracy, UX, and governance milestones",
        ],
    )
    _slide_bullets(
        prs,
        "Thank you",
        [
            "Questions and deep-dive sessions welcome",
            f"Deck date: {TODAY}",
        ],
    )

    prs.save(path)
    print(f"Wrote {path}")


def build_ppt_technical(path: str) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)

    _slide_title_only(
        prs,
        "Karnataka Schemes AI — Technical Deep Dive",
        f"Architecture, modules, stack · {TODAY}",
    )
    _slide_bullets(
        prs,
        "Repository-shaped summary",
        [
            "content/: ingest, RAG orchestration, Streamlit PoC, Flask app, PDF utilities",
            "Chroma persistent DB + sentence-transformers embeddings (default MiniLM L6 v2)",
            "Gemini via google-genai; env: GEMINI_API_KEY, optional GEMINI_CONTEXT_CHUNK_CHARS",
        ],
    )
    _slide_bullets(
        prs,
        "Ingestion pipeline",
        [
            "ingest_karnataka_schemes.py: pypdf → per-page text → scheme-block heuristics",
            "Metadata: scheme_name, page, category, tags, source path",
            "Cypher-style rules compiled to data/graph_compiled.json (no Neo4j in PoC)",
        ],
    )
    _slide_bullets(
        prs,
        "orchestrated_rag_schemes.py",
        [
            "_retrieve / _format_context_block: Chroma query + numbered context for LLM",
            "Optional _ask_gemini_for_refinements → extra _retrieve passes",
            "graph_addon_for_metadatas: inject filtered graph statements",
            "run_orchestrated_rag_result: answer + citations + debug dict",
        ],
    )
    _slide_bullets(
        prs,
        "Interfaces",
        [
            "poc_streamlit.py: run_orchestrated_rag_result, scheme picker, poc_pdf kit",
            "scheme_source_pdf.py: match section in source PDF for kit body (fallback: indexed chunks)",
            "app.py: Flask /api/ask → same orchestration",
        ],
    )
    _slide_bullets(
        prs,
        "Supporting modules",
        [
            "graph_scheme_documents.py: DEPENDS_ON → required documents",
            "user_profile.py / scheme_choice_label.py: question composition and labels",
            "poc_nearby.py: Nominatim + haversine to CSV centres",
        ],
    )
    _slide_bullets(
        prs,
        "Stack table",
        [
            "Vector: Chroma · Embeddings: SentenceTransformers · PDF: pypdf / fpdf2",
            "LLM: Gemini · UI: Streamlit, Flask · Geo: Nominatim (PoC)",
        ],
    )
    _slide_bullets(
        prs,
        "Alternative components",
        [
            "Vector: Qdrant, Milvus, pgvector · LLM: OpenAI, Anthropic, local vLLM",
            "Orchestration: LangChain / LlamaIndex if customer standardizes",
        ],
    )
    _slide_bullets(
        prs,
        "Planned guardrails (design)",
        [
            "Distance thresholds, abstention, structured output + verification pass",
            "PII redaction hooks in logs; input/output policy filters",
        ],
    )
    _slide_bullets(
        prs,
        "Planned privacy / governance (design)",
        [
            "Key rotation, RBAC, corpus versioning, audit metadata (model + index id)",
            "Vendor DPAs and regional deployment options per customer",
        ],
    )
    _slide_bullets(
        prs,
        "Limitations (PoC)",
        [
            "Not full production monitoring, auth, or multi-tenant isolation",
            "Guardrails/governance items above are roadmap, not shipped behaviour",
        ],
    )
    _slide_bullets(
        prs,
        "References",
        [
            "NIST AI RMF, ISO/IEC 42001 for enterprise risk conversations",
            "EU AI Act where applicable",
        ],
    )
    _slide_bullets(
        prs,
        "Configuration (high level)",
        [
            "GEMINI_API_KEY required; GEMINI_CONTEXT_CHUNK_CHARS caps context per chunk",
            "GEMINI_CALL_SPACING_S / skip-refinement env flags manage quota bursts",
            "Chroma path + collection + embedding model must match ingest settings",
        ],
    )
    _slide_bullets(
        prs,
        "Reproducing locally",
        [
            "pip install -r requirements-poc.txt (+ chromadb, pypdf, sentence-transformers)",
            "Run ingest on data/Karnataka Schemes.pdf → persistent Chroma DB",
            "streamlit run poc_streamlit.py or Flask app.py for /api/ask",
        ],
    )
    _slide_bullets(
        prs,
        "Q&A",
        [
            "Technical deep dive: code walkthrough, deployment sketch, load tests on request",
            f"{TODAY}",
        ],
    )

    prs.save(path)
    print(f"Wrote {path}")


def main() -> None:
    os.makedirs(OUT_DIR, exist_ok=True)
    docx_path = os.path.join(OUT_DIR, "Karnataka_Schemes_PoC_Writeup.docx")
    ppt_exec = os.path.join(OUT_DIR, "PoC_Executive_Overview_15slides.pptx")
    ppt_tech = os.path.join(OUT_DIR, "PoC_Technical_Overview_15slides.pptx")

    build_docx(docx_path)
    build_ppt_executive(ppt_exec)
    build_ppt_technical(ppt_tech)
    print(f"All files under: {OUT_DIR}")


if __name__ == "__main__":
    main()
